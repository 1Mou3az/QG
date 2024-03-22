import PyPDF2
import requests
from bs4 import BeautifulSoup
import csv
import zipfile
import json
import xml.etree.ElementTree as ET
import docx
import pptx
import rarfile
import openpyxl 
from youtube_transcript_api import YouTubeTranscriptApi
import pke
from nltk.stem import PorterStemmer
from sklearn.metrics.pairwise import cosine_similarity
from sklearn.feature_extraction.text import CountVectorizer
from sense2vec import Sense2Vec
from collections import OrderedDict
import os
from peft import PeftModel, PeftConfig
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
import re
from nltk.tokenize import word_tokenize

s2v = Sense2Vec().from_disk('s2v_old')

# Load the Lora model
config = PeftConfig.from_pretrained("mou3az/Generate-Questions")
model = AutoModelForSeq2SeqLM.from_pretrained("facebook/bart-base")
G_tokenizer = AutoTokenizer.from_pretrained("facebook/bart-base")
G_model = PeftModel.from_pretrained(model, "mou3az/Generate-Questions")


def extract_video_id(link):
    # Extract video ID from various YouTube link formats
    video_id_match = re.search(r"(?:youtube\.com\/(?:[^\/\n\s]+\/\S+\/|(?:v|e(?:mbed)?)\/|\S*?[?&]v=)|youtu\.be\/)([a-zA-Z0-9_-]{11})", link)
    if video_id_match:
        return video_id_match.group(1)
    else:
        return None
    
def get_youtube_transcripts(video_id):
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = ' '.join([line['text'] for line in transcript])
        return text
    except Exception as e:
        return None

def read_csv(file_path):
    with open(file_path, 'r', newline='') as csvfile:
        csv_reader = csv.reader(csvfile)
        csv_data = [row for row in csv_reader]
    return csv_data

def read_text(file_path):
    with open(file_path, 'r') as f:
        text_data = f.read()
    return text_data

def read_pdf(file_path):
    pdf_file = open(file_path, 'rb')
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text_data = []
    for page in pdf_reader.pages:
        text_data.append(page.extract_text())
    return text_data

def read_web_page(url):
    result = requests.get(url)
    src = result.content
    soup = BeautifulSoup(src, 'html.parser')
    text_data = ''
    for p in soup.find_all('p'):
        text_data += p.get_text() + '\n'
    return text_data

def read_docx(file_path):
    doc = docx.Document(file_path)
    text_data = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
    return text_data

def read_pptx(file_path):
    ppt = pptx.Presentation(file_path)
    text_data = ''
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_data += shape.text + '\n'
    return text_data

def read_xlsx(file_path):
    workbook = openpyxl.load_workbook(file_path)
    sheet = workbook.active
    text_data = ''
    for row in sheet.iter_rows(values_only=True):
        text_data += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
    return text_data

def read_json(file_path):
    with open(file_path, 'r') as f:
        json_data = json.load(f)
    return json_data

def read_html(file_path):
    with open(file_path, 'r') as f:
        html_data = f.read()
    return html_data

def read_xml(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()
    return ET.tostring(root, encoding='unicode')

def read_zip(file_path):
    file_contents = []
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        for file_info in zip_ref.infolist():
            with zip_ref.open(file_info) as file:
                # Call read_data to handle reading and processing the file contents
                file_data = read_data(file)
                file_contents.append(file_data)
    return file_contents

def read_rar(file_path):
    file_contents = []
    with rarfile.RarFile(file_path, 'r') as rar_ref:
        for rar_info in rar_ref.infolist():
            with rar_ref.open(rar_info) as file:
                # Call read_data to handle reading and processing the file contents
                file_data = read_data(file)
                file_contents.append(file_data)
    return file_contents

def read_data(_path=None, _link=None, text=None):
    if text is not None:
        return text
    elif _link is not None:
        if "youtube.com" in _link or "youtu.be" in _link:
            video_id = extract_video_id(_link)
            if video_id:
                try:
                    transcripts = get_youtube_transcripts(video_id)
                    if transcripts:
                        return transcripts
                    else:
                        raise ValueError("Failed to fetch transcripts from the YouTube video.")
                except Exception as e:
                    raise ValueError(f"An error occurred while fetching transcripts: {str(e)}")
            else:
                raise ValueError("Invalid YouTube link provided.")
        elif _link.startswith('https'): 
            return read_web_page(_link)

    elif _path is not None:
        # If a file path is provided, determine its type and read accordingly
        if _path.endswith('.csv'):
            return read_csv(_path)
        elif _path.endswith('.txt'):
            return read_text(_path)
        elif _path.endswith('.pdf'):
            return read_pdf(_path)
        elif _path.endswith('.docx'):
            return read_docx(_path)
        elif _path.endswith('.html'):
            return read_html(_path)
        elif _path.endswith('.pptx'):
            return read_pptx(_path)
        elif _path.endswith('.xlsx'):
            return read_xlsx(_path)
        elif _path.endswith('.json'):
            return read_json(_path)
        elif _path.endswith('.xml'):
            return read_xml(_path)
        elif _path.endswith('.zip'):
            return read_zip(_path)
        elif _path.endswith('.rar'):
            return read_rar(_path)
        else:
            raise ValueError("Unsupported type")

    else:
        # If neither file path nor text is provided, return None
        return None  # Placeholder for handling other cases

def remove_non_ascii(text):
    """Remove non-ASCII characters from list of tokenized words"""
    return text.encode('ascii','ignore').decode()

def remove_brackets_num(text):
    return re.sub(r"\*?", "", text)

def to_lowercase(text):
    return text.lower()

def replace_numbers(text):
    """Replace all interger occurrences in list of tokenized words with textual representation"""
    return re.sub(r'\d+','',text)

def remove_whitespace(text):
      return text.strip()

def remove_punctuation(text):
    punctuation= r'''!()[]{};:'"<>/?$%^&*_`~='''
    for punc in punctuation:
        text=text.replace(punc,"")
    return text

def remove_emails(text):
    return re.sub(r'[A-Za-z0-9]*@[A-Za-z]*\.?[A-Za-z0-9]*', "", text)

def text2words(text):
    return word_tokenize(text)

def normalize_text(text):
    text = remove_non_ascii(text)
    text= remove_brackets_num(text)
    text = to_lowercase(text)
    #text=replace_numbers(text)
    text= remove_whitespace(text)
    text = remove_punctuation(text)
    text= remove_emails(text)
    words = text2words(text)

    return ' '.join(words)

def calculate_similarity(sentence1, sentence2):
    # Initialize Porter Stemmer
    stemmer = PorterStemmer()

    # Tokenize and stem the sentences
    stemmed_sentence1 = ' '.join([stemmer.stem(word) for word in sentence1.split()])
    stemmed_sentence2 = ' '.join([stemmer.stem(word) for word in sentence2.split()])

    # Convert the stemmed sentences into vectors
    vectorizer = CountVectorizer().fit([stemmed_sentence1, stemmed_sentence2])
    vectorized_sentences = vectorizer.transform([stemmed_sentence1, stemmed_sentence2])

    # Calculate cosine similarity
    cosine_sim = cosine_similarity(vectorized_sentences)[0][1]

    return cosine_sim

def extract_keywords_from_text(text):
    # Initialize keyphrase extraction model, here TopicRank
    extractor = pke.unsupervised.TopicRank()

    # Load the content of the document
    extractor.load_document(input=text, language='en')

    # Keyphrase candidate selection: in the case of TopicRank: sequences of nouns
    # and adjectives (i.e., `(Noun|Adj)*`)
    extractor.candidate_selection()

    # Candidate weighting: using a random walk algorithm
    extractor.candidate_weighting()

    # N-best selection, keyphrases contains the 10 highest scored candidates
    keyphrases = extractor.get_n_best(n=20)

    # Extract keyphrases
    keywords = [keyphrase for keyphrase, score in keyphrases]

    # Calculate similarity with keywords
    unique_keywords = []

    # Handling unigrams and bigrams separately
#     unigrams = [keyphrase for keyphrase in keywords if len(keyphrase.split()) == 1]
    bigrams = [keyphrase for keyphrase in keywords if len(keyphrase.split()) == 2]

    # Add unigrams to unique_keywords directly
#     unique_keywords.extend(unigrams)

    # Filter bigrams based on similarity with existing keywords
    for keyphrase in bigrams:
        similarity = calculate_similarity(keyphrase, ' '.join(keywords))
        if similarity < 0.4:  # Adjust the similarity threshold as needed
            unique_keywords.append(keyphrase)

    return unique_keywords

def sense2vec_get_words(word, s2v):
    output = []
    word = word.lower()

    sense = s2v.get_best_sense(word)
    similarity_threshold = 0.4
    out = []
    
    if sense is not None:
        most_similar = s2v.most_similar(sense, n=20)
        for sim in most_similar:
            append_word= sim[0].split("|")[0].replace("_", " ").lower()

            # Check similarity with keyword
            similarity_keyword = calculate_similarity(word, append_word)
            #print(f"Similarity between '{word}' and '{append_word}': {similarity_keyword}")

            # Check if similarity with keyword is above the threshold
            if similarity_keyword >= similarity_threshold:
                continue
            
            # Check similarity with existing distractors
            similarity_to_existing = [calculate_similarity(append_word, existing_distractor) for existing_distractor in output]

            # Check if similarity with any existing distractor is above the threshold
            if any(similarity >= similarity_threshold for similarity in similarity_to_existing):
                continue
            
            # If the conditions are met, append the word to the list of output
            output.append(append_word.title())

        out = list(OrderedDict.fromkeys(output))
    return out[:3]

def generate_questions(context, answer):
    device = next(G_model.parameters()).device
    input_text = f"Given the context '{context}' and the answer '{answer}', what question can be asked?"
    encoding = G_tokenizer.encode_plus(input_text, padding=True, return_tensors="pt").to(device)

    output_tokens = G_model.generate(
        **encoding,
        early_stopping=True,
        do_sample= True,
        num_beams=5,
        num_return_sequences=1,
        no_repeat_ngram_size=2,
        max_length=200,
        temperature=0.6,
        top_p=0.95,
        repetition_penalty=1.2
    )
    question = G_tokenizer.decode(output_tokens[0], skip_special_tokens=True).replace("question:", "").strip()
    return question

def process_and_generate_questions(_path=None, _link=None, text=None):
    N_text_file = ""
    keyword_question_distractors = []

    try:
        # Input Validation
        if (_path is None and _link is None and text is None) or (_path and not os.path.exists(_path)):
            raise ValueError("Invalid input: Please provide a valid file path or link or text.")

        # Read the data from the file
        try:
            example = read_data(_path, _link, text)
        except ValueError as e:
            return None, str(e)
        
        # Check if example is not None and not empty
        if example is None or not example.strip():
            raise ValueError("Empty or invalid input: Please provide a valid file or non-empty text or page.")

        # Preprocessing
        processed_text=normalize_text(example)

        # Tokenization
        tokenized_text = G_tokenizer.encode_plus(processed_text)
        if len(tokenized_text["input_ids"]) > 1024:
            raise ValueError("Tokenized text exceeds maximum length (1024 tokens).")

        # Normalization
        N_text_file = processed_text

        if not N_text_file:
            raise ValueError("Normalization failed: Unable to process text.")

        # Extracting Keywords and Generating Questions
        keywords = extract_keywords_from_text(N_text_file)
        if not keywords:
            return None, "No keywords generated."
        else:
            # Process keywords and generate questions
            for keyword in keywords:
                try:
                    current_distractors = sense2vec_get_words(keyword, s2v)
                    if current_distractors:
                        question = generate_questions(N_text_file, keyword)
                        keyword_question_distractors.append((keyword, current_distractors, question))
                except Exception as e:
                    raise ValueError(f"An error occurred while processing keyword '{keyword}': {e}")

        if not keyword_question_distractors:
            return None, "No questions could be generated due to the lack of keywords and distractors."

    except ValueError as ve:
        return None, str(ve)

    except Exception as e:
        return None, f"An unexpected error occurred: {str(e)}"

    return {'N_text_file': N_text_file, 'keyword_question_distractors': keyword_question_distractors}, None